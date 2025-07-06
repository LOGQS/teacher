#!/usr/bin/env python3
"""
PowerPoint Presentation Builder using python-pptx
Builds complete PowerPoint presentations from slide content specifications.
"""

import logging
from pathlib import Path
from typing import Dict, List, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

class PresentationBuilder:
    """Builds PowerPoint presentations using python-pptx"""
    
    def __init__(self, file_manager=None):
        """Initialize the presentation builder"""
        self.file_manager = file_manager
        # Presentation files will be saved in session directories
        
        # Default theme settings
        self.themes = {
            'default': {
                'background_color': RGBColor(32, 32, 32),  # Dark theme
                'title_color': RGBColor(255, 255, 255),
                'text_color': RGBColor(230, 230, 230),
                'accent_color': RGBColor(0, 162, 255),  # iOS blue
                'title_font': 'SF Pro Display',
                'body_font': 'SF Pro Text',
                'fallback_title_font': 'Segoe UI',
                'fallback_body_font': 'Segoe UI'
            },
            'light': {
                'background_color': RGBColor(248, 248, 248),
                'title_color': RGBColor(32, 32, 32),
                'text_color': RGBColor(64, 64, 64),
                'accent_color': RGBColor(0, 122, 255),
                'title_font': 'SF Pro Display',
                'body_font': 'SF Pro Text',
                'fallback_title_font': 'Segoe UI',
                'fallback_body_font': 'Segoe UI'
            }
        }
    
    def _get_presentation_path(self, session_id: str) -> Path:
        """Get presentation file path for a specific session"""
        if self.file_manager:
            session_dir = self.file_manager.get_session_dir(session_id)
            return session_dir / "presentation.pptx"
        else:
            # Fallback to old structure if no file_manager
            output_dir = Path('data/presentations')
            output_dir.mkdir(parents=True, exist_ok=True)
            filename = self._generate_filename("presentation")
            return output_dir / filename
    
    def build_presentation(self, 
                          slides_content: List[Dict[str, Any]],
                          course_title: str,
                          session_id: str = None,
                          theme: str = 'default') -> str:
        """
        Build complete PowerPoint presentation from slide content
        
        Args:
            slides_content: List of detailed slide content dictionaries
            course_title: Main course title
            session_id: Session identifier for organizing files
            theme: Theme name for styling
            
        Returns:
            Path to the generated presentation file
        """
        try:
            # Create new presentation
            prs = Presentation()
            
            # Get theme settings
            theme_config = self.themes.get(theme, self.themes['default'])
            
            # Set presentation properties
            prs.core_properties.title = course_title
            prs.core_properties.author = "AI-Powered Educational System"
            prs.core_properties.subject = "Educational Presentation"
            
            logger.info(f"Building presentation '{course_title}' with {len(slides_content)} slides")
            
            # Enhance slide layouts for better presentation quality
            enhanced_slides = self.enhance_slide_layout(slides_content)
            
            # Process each slide
            for slide_idx, slide_data in enumerate(enhanced_slides):
                try:
                    self._build_slide(prs, slide_data, theme_config, slide_idx + 1)
                    logger.debug(f"Built slide {slide_idx + 1}: {slide_data.get('title', 'Untitled')}")
                except Exception as e:
                    logger.error(f"Error building slide {slide_idx + 1}: {str(e)}")
                    # Create a basic error slide
                    self._create_error_slide(prs, slide_idx + 1, str(e))
            
            # Save presentation using session-based path
            if session_id:
                file_path = self._get_presentation_path(session_id)
            else:
                # Fallback to old behavior
                filename = self._generate_filename(course_title)
                output_dir = Path('data/presentations')
                output_dir.mkdir(parents=True, exist_ok=True)
                file_path = output_dir / filename
            
            prs.save(str(file_path))
            
            logger.info(f"Successfully saved presentation: {file_path}")
            return str(file_path)
            
        except Exception as e:
            logger.error(f"Error building presentation: {str(e)}")
            raise
    
    def enhance_slide_layout(self, slides_content: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """
        Enhance slide layout specifications for better presentation quality
        
        Args:
            slides_content: Original slide content
            
        Returns:
            Enhanced slide content with improved layouts
        """
        try:
            enhanced_slides = []
            
            for slide_data in slides_content:
                enhanced_slide = slide_data.copy()
                layout = slide_data.get('layout', {})
                
                # If no explicit layout, create a default one
                if not layout.get('elements'):
                    enhanced_slide['layout'] = self._create_default_layout(slide_data)
                else:
                    # Enhance existing layout
                    enhanced_slide['layout'] = self._enhance_existing_layout(layout, slide_data)
                
                enhanced_slides.append(enhanced_slide)
            
            return enhanced_slides
            
        except Exception as e:
            logger.error(f"Error enhancing slide layouts: {str(e)}")
            return slides_content  # Return original if enhancement fails
    
    def _create_default_layout(self, slide_data: Dict[str, Any]) -> Dict[str, Any]:
        """Create a default layout for slide with no explicit layout"""
        title = slide_data.get('title', '')
        transcript = slide_data.get('transcript', '')
        images = slide_data.get('images', [])
        
        layout = {
            'slide_type': 'content_slide',
            'background_color': None,  # Use theme default
            'elements': []
        }
        
        # Add title element
        if title:
            layout['elements'].append({
                'type': 'title',
                'content': title,
                'position': {'x': 0.5, 'y': 0.5, 'width': 9, 'height': 1},
                'formatting': {'size': 28, 'bold': True, 'alignment': 'center'}
            })
        
        # Add content based on transcript length
        if transcript:
            # Split transcript into bullet points if it's long
            if len(transcript) > 200:
                # Try to split into logical sections
                sentences = transcript.split('. ')
                if len(sentences) > 3:
                    # Create bullet points from first few sentences
                    bullet_points = [f"{sentence.strip()}." for sentence in sentences[:4] if sentence.strip()]
                    layout['elements'].append({
                        'type': 'textbox',
                        'content': bullet_points,
                        'position': {'x': 1, 'y': 2, 'width': 8, 'height': 4},
                        'formatting': {'size': 18, 'alignment': 'left'}
                    })
                else:
                    # Add as single text block
                    layout['elements'].append({
                        'type': 'textbox',
                        'content': transcript,
                        'position': {'x': 1, 'y': 2, 'width': 8, 'height': 3},
                        'formatting': {'size': 16, 'alignment': 'left'}
                    })
            else:
                # Short transcript, add as is
                layout['elements'].append({
                    'type': 'textbox',
                    'content': transcript,
                    'position': {'x': 1, 'y': 2, 'width': 8, 'height': 2},
                    'formatting': {'size': 18, 'alignment': 'left'}
                })
        
        # Adjust layout if images are present
        if images:
            # Modify text positioning to accommodate images
            for element in layout['elements']:
                if element['type'] == 'textbox':
                    if len(images) <= 2:
                        # Images on the right, text on the left
                        element['position']['width'] = 5
                    else:
                        # Images at bottom, text at top
                        element['position']['height'] = 2.5
        
        return layout
    
    def _enhance_existing_layout(self, layout: Dict[str, Any], slide_data: Dict[str, Any]) -> Dict[str, Any]:
        """Enhance existing layout with better positioning and formatting"""
        enhanced_layout = layout.copy()
        
        # Ensure all elements have proper positioning
        for element in enhanced_layout.get('elements', []):
            if 'position' not in element:
                element['position'] = self._get_default_position(element.get('type', 'textbox'))
            
            if 'formatting' not in element:
                element['formatting'] = self._get_default_formatting(element.get('type', 'textbox'))
        
        return enhanced_layout
    
    def _get_default_position(self, element_type: str) -> Dict[str, float]:
        """Get default position for element type"""
        positions = {
            'title': {'x': 0.5, 'y': 0.5, 'width': 9, 'height': 1},
            'textbox': {'x': 1, 'y': 2, 'width': 8, 'height': 3},
            'shape': {'x': 4, 'y': 3, 'width': 2, 'height': 1},
            'image': {'x': 6, 'y': 2, 'width': 3, 'height': 2}
        }
        return positions.get(element_type, positions['textbox'])
    
    def _get_default_formatting(self, element_type: str) -> Dict[str, Any]:
        """Get default formatting for element type"""
        formatting = {
            'title': {'size': 28, 'bold': True, 'alignment': 'center'},
            'textbox': {'size': 18, 'alignment': 'left'},
            'shape': {'fill_color': '#0080FF'},
        }
        return formatting.get(element_type, formatting['textbox'])
    
    def _build_slide(self, 
                    prs: Presentation, 
                    slide_data: Dict[str, Any], 
                    theme_config: Dict[str, Any],
                    slide_number: int):
        """Build a single slide from slide data"""
        try:
            layout = slide_data.get('layout', {})
            slide_type = layout.get('slide_type', 'content_slide')
            
            # Choose appropriate slide layout
            if slide_type == 'title_slide' or slide_number == 1:
                slide_layout = prs.slide_layouts[0]  # Title slide
            elif slide_type == 'section_header':
                slide_layout = prs.slide_layouts[2]  # Section header
            else:
                slide_layout = prs.slide_layouts[6]  # Blank layout for custom design
            
            # Add slide
            slide = prs.slides.add_slide(slide_layout)
            
            # Set background
            self._set_slide_background(slide, layout, theme_config)
            
            # Add slide elements
            self._add_slide_elements(slide, slide_data, theme_config)
            
            # Add images
            self._add_slide_images(slide, slide_data)
            
            # Add speaker notes
            self._add_speaker_notes(slide, slide_data.get('transcript', ''))
            
        except Exception as e:
            logger.error(f"Error building individual slide: {str(e)}")
            raise
    
    def _set_slide_background(self, slide, layout: Dict[str, Any], theme_config: Dict[str, Any]):
        """Set slide background color"""
        try:
            background_color = layout.get('background_color')
            if background_color:
                # Parse color if it's a string
                if isinstance(background_color, str):
                    if background_color.startswith('#'):
                        # Convert hex to RGB
                        hex_color = background_color[1:]
                        rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                        color = RGBColor(*rgb)
                    else:
                        color = theme_config['background_color']
                else:
                    color = background_color
            else:
                color = theme_config['background_color']
            
            # Set background fill
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = color
            
        except Exception as e:
            logger.warning(f"Error setting background: {str(e)}")
            # Use default background
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = theme_config['background_color']
    
    def _add_slide_elements(self, slide, slide_data: Dict[str, Any], theme_config: Dict[str, Any]):
        """Add textboxes and other elements to the slide"""
        try:
            layout = slide_data.get('layout', {})
            elements = layout.get('elements', [])
            
            # Add title if not in elements
            title = slide_data.get('title', '')
            if title and not any(e.get('type') == 'title' for e in elements):
                self._add_title_element(slide, title, theme_config)
            
            # Add each element
            for element in elements:
                element_type = element.get('type', 'textbox')
                
                if element_type == 'textbox':
                    self._add_textbox_element(slide, element, theme_config)
                elif element_type == 'title':
                    self._add_title_element(slide, element.get('content', ''), theme_config, element)
                elif element_type == 'shape':
                    self._add_shape_element(slide, element, theme_config)
                
        except Exception as e:
            logger.warning(f"Error adding slide elements: {str(e)}")
    
    def _add_title_element(self, slide, title: str, theme_config: Dict[str, Any], element: Dict[str, Any] = None):
        """Add title element to slide"""
        try:
            if element and 'position' in element:
                pos = element['position']
                left = Inches(pos.get('x', 0.5))
                top = Inches(pos.get('y', 0.5))
                width = Inches(pos.get('width', 9))
                height = Inches(pos.get('height', 1))
            else:
                # Default title position
                left = Inches(0.5)
                top = Inches(0.5)
                width = Inches(9)
                height = Inches(1)
            
            # Add title textbox
            title_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = title_box.text_frame
            text_frame.clear()
            
            p = text_frame.paragraphs[0]
            p.text = title
            p.alignment = PP_ALIGN.CENTER
            
            # Format title
            run = p.runs[0]
            font = run.font
            font.name = theme_config.get('title_font', theme_config.get('fallback_title_font', 'Arial'))
            font.size = Pt(element.get('formatting', {}).get('size', 32) if element else 32)
            font.color.rgb = theme_config['title_color']
            font.bold = True
            
        except Exception as e:
            logger.warning(f"Error adding title: {str(e)}")
    
    def _add_textbox_element(self, slide, element: Dict[str, Any], theme_config: Dict[str, Any]):
        """Add textbox element to slide"""
        try:
            position = element.get('position', {})
            content = element.get('content', '')
            formatting = element.get('formatting', {})
            
            # Position
            left = Inches(position.get('x', 1))
            top = Inches(position.get('y', 2))
            width = Inches(position.get('width', 8))
            height = Inches(position.get('height', 1))
            
            # Add textbox
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.clear()
            text_frame.word_wrap = True
            
            # Add content
            if isinstance(content, list):
                # Handle bullet points
                for i, item in enumerate(content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = str(item)
                    p.level = 0
            else:
                p = text_frame.paragraphs[0]
                p.text = str(content)
            
            # Apply formatting
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = self._get_alignment(formatting.get('alignment', 'left'))
                
                for run in paragraph.runs:
                    font = run.font
                    font.name = theme_config.get('body_font', theme_config.get('fallback_body_font', 'Arial'))
                    font.size = Pt(formatting.get('size', 18))
                    font.color.rgb = theme_config['text_color']
                    
                    if formatting.get('bold'):
                        font.bold = True
                    if formatting.get('italic'):
                        font.italic = True
            
        except Exception as e:
            logger.warning(f"Error adding textbox: {str(e)}")
    
    def _add_shape_element(self, slide, element: Dict[str, Any], theme_config: Dict[str, Any]):
        """Add shape element to slide"""
        try:
            position = element.get('position', {})
            shape_type = element.get('shape_type', 'rectangle')
            
            # Position
            left = Inches(position.get('x', 1))
            top = Inches(position.get('y', 2))
            width = Inches(position.get('width', 2))
            height = Inches(position.get('height', 1))
            
            # Map shape types
            shape_map = {
                'rectangle': MSO_SHAPE.RECTANGLE,
                'oval': MSO_SHAPE.OVAL,
                'triangle': MSO_SHAPE.ISOSCELES_TRIANGLE,  # Use specific triangle type
                'diamond': MSO_SHAPE.DIAMOND,
                'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
                'arrow': MSO_SHAPE.RIGHT_ARROW,
                'star': MSO_SHAPE.STAR_5_POINT
            }
            
            shape_type_enum = shape_map.get(shape_type.lower(), MSO_SHAPE.RECTANGLE)
            
            # Add shape with error handling
            try:
                shape = slide.shapes.add_shape(shape_type_enum, left, top, width, height)
            except Exception as shape_error:
                logger.warning(f"Failed to add shape {shape_type}, using rectangle: {str(shape_error)}")
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            
            # Format shape
            if 'fill_color' in element:
                shape.fill.solid()
                color = self._parse_color(element['fill_color'], theme_config['accent_color'])
                shape.fill.fore_color.rgb = color
            
            if 'line_color' in element:
                color = self._parse_color(element['line_color'], theme_config['text_color'])
                shape.line.color.rgb = color
            
        except Exception as e:
            logger.warning(f"Error adding shape: {str(e)}")
    
    def _add_slide_images(self, slide, slide_data: Dict[str, Any]):
        """Add images to the slide"""
        try:
            images = slide_data.get('images', [])
            
            for image_data in images:
                file_path = image_data.get('file_path')
                if not file_path or not Path(file_path).exists():
                    logger.warning(f"Image file not found: {file_path}")
                    continue
                
                position = image_data.get('position', {})
                
                # Position
                left = Inches(position.get('x', 1))
                top = Inches(position.get('y', 3))
                width = Inches(position.get('width', 3))
                height = Inches(position.get('height', 2))
                
                try:
                    # Add image
                    pic = slide.shapes.add_picture(file_path, left, top, width, height)
                    
                    # Add caption if provided
                    caption = image_data.get('caption')
                    if caption:
                        caption_top = top + height + Inches(0.1)
                        caption_box = slide.shapes.add_textbox(left, caption_top, width, Inches(0.3))
                        caption_frame = caption_box.text_frame
                        caption_frame.text = caption
                        
                        # Format caption
                        p = caption_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        run = p.runs[0]
                        run.font.size = Pt(12)
                        run.font.italic = True
                    
                except Exception as e:
                    logger.warning(f"Error adding image {file_path}: {str(e)}")
                    
        except Exception as e:
            logger.warning(f"Error adding images: {str(e)}")
    
    def _add_speaker_notes(self, slide, transcript: str):
        """Add speaker notes to the slide"""
        try:
            if transcript:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = transcript
                
        except Exception as e:
            logger.warning(f"Error adding speaker notes: {str(e)}")
    
    def _create_error_slide(self, prs: Presentation, slide_number: int, error_message: str):
        """Create a basic error slide"""
        try:
            slide_layout = prs.slide_layouts[6]  # Blank
            slide = prs.slides.add_slide(slide_layout)
            
            # Add error message
            textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
            text_frame = textbox.text_frame
            text_frame.text = f"Error building slide {slide_number}:\n{error_message}"
            
            # Format error text
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(255, 0, 0)
            
        except Exception as e:
            logger.error(f"Error creating error slide: {str(e)}")
    
    def _get_alignment(self, alignment_str: str):
        """Convert alignment string to PowerPoint alignment"""
        alignment_map = {
            'left': PP_ALIGN.LEFT,
            'center': PP_ALIGN.CENTER,
            'right': PP_ALIGN.RIGHT,
            'justify': PP_ALIGN.JUSTIFY
        }
        return alignment_map.get(alignment_str.lower(), PP_ALIGN.LEFT)
    
    def _parse_color(self, color_spec: Any, default_color: RGBColor) -> RGBColor:
        """Parse color specification"""
        try:
            if isinstance(color_spec, str):
                if color_spec.startswith('#'):
                    hex_color = color_spec[1:]
                    rgb = tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))
                    return RGBColor(*rgb)
            elif isinstance(color_spec, dict) and 'rgb' in color_spec:
                rgb = color_spec['rgb']
                return RGBColor(rgb[0], rgb[1], rgb[2])
            elif isinstance(color_spec, (list, tuple)) and len(color_spec) == 3:
                return RGBColor(*color_spec)
            
            return default_color
            
        except Exception:
            return default_color
    
    def _generate_filename(self, course_title: str) -> str:
        """Generate filename for the presentation"""
        try:
            # Clean title for filename
            import re
            from datetime import datetime
            
            clean_title = re.sub(r'[^\w\s-]', '', course_title)
            clean_title = re.sub(r'[-\s]+', '_', clean_title)
            clean_title = clean_title.strip('_')[:50]  # Limit length
            
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            
            return f"{clean_title}_{timestamp}.pptx"
            
        except Exception:
            from datetime import datetime
            timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
            return f"presentation_{timestamp}.pptx"